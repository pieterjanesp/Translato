from app.handlers.base import BaseDocumentHandler
from pptx import Presentation

class PPTXHandler(BaseDocumentHandler):
    """Handler for PowerPoint files."""

    def extract_text(self, file_path: str) -> list[str]:
        """Extract text from PPTX file."""
        presentation = Presentation(file_path)
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.strip():
                        texts.append(shape.text)
        return texts

    def replace_text(self, file_path: str, translations: dict[str, str], output_path: str) -> None:
        """Replace text in PPTX file with translations."""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for original,translated in translations.items():
                                if original in run.text:
                                    run.text = run.text.replace(original, translated)
        presentation.save(output_path)
        

